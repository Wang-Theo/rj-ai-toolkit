"""
PowerPoint 文件解析模块
解析 PPTX 文件，将内容转换为 Markdown 格式，表格保留为 HTML 格式
"""

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pathlib import Path
from typing import List, Union, Optional
from io import BytesIO
import re
from PIL import Image
import base64
import tempfile
from .base_parser import BaseParser


class PPTXParser(BaseParser):
    """PowerPoint 文件解析器"""
    
    def __init__(
        self, 
        use_ocr: bool = False, 
        ocr_func = None,
        save_images: bool = True,
        images_dir: Optional[Union[str, Path]] = None,
        image_dpi: int = 300
    ):
        """
        初始化 PPTX 解析器
        
        Args:
            use_ocr: 是否使用OCR功能，默认False
            ocr_func: OCR处理函数，接收图片路径返回文本。签名: func(image_path: str) -> str
            save_images: 是否保存图片，默认True
            images_dir: 图片保存目录，默认为None（与PPTX文件同目录创建images文件夹）
            image_dpi: 图片 DPI，默认300
        """
        super().__init__(use_ocr, ocr_func)
        self.save_images = save_images
        self.images_dir = Path(images_dir) if images_dir else None
        self.image_counter = 0
        self.image_dpi = image_dpi
    
    def parse_file(self, file_path: Union[str, Path]) -> str:
        """
        解析 PPTX 文件为 Markdown 格式文本
        
        Args:
            file_path: PPTX 文件路径
            
        Returns:
            Markdown 格式的文本内容
        """
        file_path = Path(file_path)
        
        # 如果未指定图片目录，使用PPTX文件同目录下的images文件夹
        if self.images_dir is None:
            images_base_dir = file_path.parent / f"{file_path.stem}_images"
        else:
            images_base_dir = self.images_dir
        
        # 重置图片计数器
        self.image_counter = 0
        
        # 解析 PPTX 文件
        prs = Presentation(str(file_path))
        markdown_content = []
        
        # 添加文档标题
        markdown_content.append(f"# {file_path.stem}\n")
        markdown_content.append("---\n")
        
        # 遍历每一页幻灯片
        for slide_idx, slide in enumerate(prs.slides, 1):
            slide_content = self._parse_slide(slide, slide_idx, images_base_dir)
            if slide_content:
                markdown_content.append(slide_content)
                markdown_content.append("\n---\n")
        
        return "\n".join(markdown_content)
    
    def get_supported_extensions(self) -> List[str]:
        """
        获取支持的文件扩展名
        
        Returns:
            支持的扩展名列表
        """
        return ['.pptx']
    
    def _parse_slide(self, slide, slide_number: int, images_base_dir: Path) -> str:
        """
        解析单个幻灯片
        
        Args:
            slide: 幻灯片对象
            slide_number: 幻灯片编号
            images_base_dir: 图片保存基础目录
            
        Returns:
            解析后的 Markdown 文本
        """
        content = []
        content.append(f"\n## Slide {slide_number}\n")
        
        # 尝试获取幻灯片标题
        if hasattr(slide, 'shapes') and hasattr(slide.shapes, 'title'):
            if slide.shapes.title and hasattr(slide.shapes.title, 'text'):
                title = slide.shapes.title.text.strip()
                if title:
                    content.append(f"### {title}\n")
        
        # 按照位置排序所有形状
        shapes_with_position = []
        for shape in slide.shapes:
            try:
                top = shape.top
                left = shape.left
                shapes_with_position.append((top, left, shape))
            except:
                shapes_with_position.append((float('inf'), float('inf'), shape))
        
        shapes_with_position.sort(key=lambda x: (x[0], x[1]))
        
        # 解析所有形状
        for _, _, shape in shapes_with_position:
            shape_content = self._parse_shape(shape, images_base_dir)
            if shape_content:
                content.append(shape_content)
        
        return "\n".join(content)
    
    def _parse_shape(self, shape, images_base_dir: Path) -> str:
        """
        解析形状对象，根据类型分发到对应的处理函数
        
        Args:
            shape: 形状对象
            images_base_dir: 图片保存基础目录
            
        Returns:
            解析后的内容
        """
        # 处理表格
        if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
            return self._parse_table(shape)
        
        # 处理图片
        elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            return self._parse_picture(shape, images_base_dir)
        
        # 处理文本框和其他包含文本的形状
        elif hasattr(shape, "text") and shape.text.strip():
            return self._parse_text_shape(shape)
        
        # 处理组合形状
        elif shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            return self._parse_group_shape(shape, images_base_dir)
        
        return ""
    
    def _parse_table(self, shape) -> str:
        """
        解析表格为 HTML 格式
        
        Args:
            shape: 表格形状对象
            
        Returns:
            HTML 格式的表格
        """
        if not hasattr(shape, 'table'):
            return ""
        
        table = shape.table
        html = ['<table border="1" cellpadding="5" cellspacing="0">']
        
        for row_idx, row in enumerate(table.rows):
            html.append('  <tr>')
            
            for cell_idx, cell in enumerate(row.cells):
                tag = 'th' if row_idx == 0 else 'td'
                
                cell_text = cell.text.strip()
                cell_text = cell_text.replace('&', '&amp;').replace('<', '&lt;').replace('>', '&gt;')
                cell_text = cell_text.replace('"', '&quot;').replace("'", '&#39;')
                
                try:
                    if cell_idx > 0 and cell == row.cells[cell_idx - 1]:
                        continue
                except:
                    pass
                
                html.append(f'    <{tag}>{cell_text}</{tag}>')
            
            html.append('  </tr>')
        
        html.append('</table>\n')
        return '\n'.join(html)
    
    def _parse_picture(self, shape, images_base_dir: Path) -> str:
        """
        解析图片形状，保存图片并返回Markdown引用
        
        Args:
            shape: 图片形状对象
            images_base_dir: 图片保存基础目录
            
        Returns:
            Markdown格式的图片引用（如果启用OCR，包含OCR识别文本）
        """
        if not hasattr(shape, 'image'):
            return ""
        
        try:
            image = shape.image
            image_bytes = image.blob
            
            # 生成图片文件名
            self.image_counter += 1
            ext = image.ext if hasattr(image, 'ext') else 'png'
            image_filename = f"image_{self.image_counter:03d}.{ext}"
            
            if self.save_images:
                # 保存图片
                images_base_dir.mkdir(parents=True, exist_ok=True)
                image_path = images_base_dir / image_filename
                
                # 使用PIL转换图片格式（确保兼容性）
                try:
                    img = Image.open(BytesIO(image_bytes))
                    
                    # 转换为RGB模式
                    if img.mode in ('RGBA', 'LA', 'P'):
                        background = Image.new('RGB', img.size, (255, 255, 255))
                        if img.mode == 'P':
                            img = img.convert('RGBA')
                        background.paste(img, mask=img.split()[-1] if img.mode in ('RGBA', 'LA') else None)
                        img = background
                    elif img.mode != 'RGB':
                        img = img.convert('RGB')
                    
                    # 保存为 PNG 格式
                    image_path = images_base_dir / f"image_{self.image_counter:03d}.png"
                    img.save(image_path, 'PNG', dpi=(self.image_dpi, self.image_dpi))
                    
                    # 返回相对路径
                    rel_path = f"{images_base_dir.name}/image_{self.image_counter:03d}.png"
                    result = f"\n![Image {self.image_counter}]({rel_path})\n"
                    
                    # 如果启用OCR，添加OCR识别内容
                    if self.use_ocr:
                        ocr_text = self._call_ocr(str(image_path))
                        if ocr_text:
                            result += f"\n**[OCR Content]:**\n\n{ocr_text}\n"
                    
                    return result
                except Exception:
                    # 如果PIL转换失败，直接保存原始字节
                    image_path.write_bytes(image_bytes)
                    rel_path = f"{images_base_dir.name}/{image_filename}"
                    result = f"\n![Image {self.image_counter}]({rel_path})\n"
                    
                    # 如果启用OCR，尝试处理
                    if self.use_ocr:
                        ocr_text = self._call_ocr(str(image_path))
                        if ocr_text:
                            result += f"\n**[OCR Content]:**\n\n{ocr_text}\n"
                    
                    return result
            else:
                # 不保存，只标记
                return f"\n[Image {self.image_counter}: {image_filename}]\n"
                
        except Exception:
            return ""
    
    def _parse_text_shape(self, shape) -> str:
        """
        解析文本形状，自动识别列表并格式化
        
        Args:
            shape: 文本形状对象
            
        Returns:
            Markdown 格式的文本或列表
        """
        # 跳过标题
        if hasattr(shape, "name") and "Title" in shape.name:
            return ""
        
        text = shape.text.strip()
        lines = text.split('\n')
        
        # 检查是否为列表项
        list_pattern = re.compile(r'^[\s]*[•●○■□▪▫◦‣⁃-]\s+|^\s*\d+[\.\)]\s+')
        list_count = sum(1 for line in lines if list_pattern.match(line))
        
        if len(lines) >= 2 and list_count >= 1:
            formatted_lines = []
            for line in lines:
                line = line.strip()
                if not line:
                    continue
                line = re.sub(r'^[•●○■□▪▫◦‣⁃]\s+', '- ', line)
                line = re.sub(r'^(\d+)[\.\)]\s+', r'\1. ', line)
                if not line.startswith(('-', '*', '+')) and not re.match(r'^\d+\.\s+', line):
                    line = f"- {line}"
                formatted_lines.append(line)
            return '\n'.join(formatted_lines) + '\n'
        else:
            return f"{text}\n"
    
    def _parse_group_shape(self, shape, images_base_dir: Path) -> str:
        """
        解析组合形状，递归处理其中的子形状
        
        Args:
            shape: 组合形状对象
            images_base_dir: 图片保存基础目录
            
        Returns:
            所有子形状解析后的内容
        """
        content = []
        for sub_shape in shape.shapes:
            sub_content = self._parse_shape(sub_shape, images_base_dir)
            if sub_content:
                content.append(sub_content)
        return "\n".join(content)
